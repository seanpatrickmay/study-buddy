"""Utilities for converting various study materials into markdown."""
from __future__ import annotations

import tempfile
from pathlib import Path
from typing import Iterable, List

import markdownify
import pptx
import mammoth

from app.services.models import DocumentBundle


class StudyConverter:
    """Convert supported study inputs (pdf/html/pptx) into markdown bundles."""

    def __init__(self, pdf_converter, html_converter) -> None:
        self._pdf_converter = pdf_converter
        self._html_converter = html_converter

    def convert_pdf(self, path: Path) -> DocumentBundle:
        markdown = self._pdf_converter(path)
        return DocumentBundle(source_path=path, display_name=path.name, markdown=markdown)

    def convert_powerpoint(self, path: Path) -> DocumentBundle:
        doc = pptx.Presentation(str(path))
        parts: List[str] = [f"# {path.stem}"]
        for slide_index, slide in enumerate(doc.slides, start=1):
            parts.append(f"\n## Slide {slide_index}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        markdown = "\n\n".join(parts)
        return DocumentBundle(source_path=path, display_name=path.name, markdown=markdown)

    def convert_html(self, path: Path) -> DocumentBundle:
        html_content = path.read_text(encoding="utf-8")
        markdown = markdownify.markdownify(html_content, heading_style="ATX")
        return DocumentBundle(source_path=path, display_name=path.name, markdown=markdown)

    def convert_docx(self, path: Path) -> DocumentBundle:
        with path.open("rb") as docx_file:
            result = mammoth.convert_to_markdown(docx_file)
        markdown = result.value
        return DocumentBundle(source_path=path, display_name=path.name, markdown=markdown)

    def convert(self, path: Path) -> DocumentBundle:
        suffix = path.suffix.lower()
        if suffix == ".pdf":
            return self.convert_pdf(path)
        if suffix in {".ppt", ".pptx"}:
            return self.convert_powerpoint(path)
        if suffix in {".html", ".htm"}:
            return self.convert_html(path)
        if suffix == ".docx":
            return self.convert_docx(path)
        raise ValueError(f"Unsupported file type: {suffix}")
